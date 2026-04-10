# -*- coding: utf-8 -*-
"""
SDACS 발표자료 v7 - PPT 생성 스크립트
- 16:9 와이드 슬라이드 16장
- 맑은 고딕 한글 폰트, 블루/화이트 테마
- v5_images/ 폴더의 15개 그림(g0–g14) 재사용
"""
import io
import os
import sys
from datetime import datetime

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


# ── 경로 설정 ────────────────────────────────────────
REPO_ROOT = r"C:\Users\sun47\Desktop\swarm-drone-atc\.claude\worktrees\stupefied-hugle"
IMG_DIR = os.path.join(REPO_ROOT, "v5_images")
OUT_DIR = os.path.join(REPO_ROOT, "docs", "presentation")
OUT_PATH = os.path.join(OUT_DIR, "SDACS_Presentation_v7.pptx")

os.makedirs(OUT_DIR, exist_ok=True)

# ── 이미지 매핑 (v7 보고서와 동일) ──────────────────────
IMG = {
    "g0":  "c6f3b3d148e074da06cbee1afedfc120699640e4.png",
    "g1":  "78cb17427a601bb078076de230d1279ef505df5a.png",
    "g2":  "e0271e0a50ac6bc78f0398f2c240f4e573bbe062.png",
    "g3":  "df06668eb85513a5ac7486fabc179561c38c8647.png",
    "g4":  "80746b62fb86c64a171c54eec1e860c582be58ce.png",
    "g5":  "cfa38868f013731afcd63c1435b362661dee6a83.png",
    "g6":  "ce7ff024213fe4ac9b79445fa88b3974a8322ea9.png",
    "g7":  "4894b105349392aeef769d0f24ad5b2d49e725a3.png",
    "g8":  "152b2021b23eb93a4da13a6a82826edb8f0e7e2c.png",
    "g9":  "f65473721ace1a8dabd5f3081c045a4e62b3dcfc.png",
    "g10": "dcd74a92f4ea60cab7a9d3640ebfa33963603eb3.png",
    "g11": "d0b4bce5521f952cd294984b8aa919b2eaaea0ea.png",
    "g12": "ac244e51158df836c161d6734de59149585a941d.png",
    "g13": "ca3c9b5e10bd3cd0f142ee2967693bf940d40e81.png",
    "g14": "7f6dcc79e78a06e547906f2d89ba69fc67bf5ad5.png",
}


# ── 테마 ────────────────────────────────────────────
ACCENT_BLUE = RGBColor(0x00, 0x5A, 0xA0)
DARK_BLUE = RGBColor(0x0A, 0x2B, 0x4F)
LIGHT_BLUE = RGBColor(0xE8, 0xF1, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)

TITLE_FONT = "맑은 고딕"
BODY_FONT = "맑은 고딕"
TITLE_PT = 28
BODY_PT = 18
SUBTITLE_PT = 20
FOOTER_PT = 10

# 16:9 슬라이드 크기
SLIDE_W_IN = 13.33
SLIDE_H_IN = 7.5


# ── 프레젠테이션 초기화 ───────────────────────────────
prs = Presentation()
prs.slide_width = Inches(SLIDE_W_IN)
prs.slide_height = Inches(SLIDE_H_IN)

BLANK_LAYOUT = prs.slide_layouts[6]  # Blank layout


# ── 헬퍼 함수 ────────────────────────────────────────
def _set_cjk_font(run, font_name: str) -> None:
    """run의 한글(동아시아) 폰트를 명시적으로 지정한다."""
    rPr = run._r.get_or_add_rPr()
    # 기존 eastAsia 노드 제거 후 추가
    for ea in rPr.findall(qn("a:ea")):
        rPr.remove(ea)
    ea = rPr.makeelement(qn("a:ea"), {"typeface": font_name})
    rPr.append(ea)


def _style_run(run, *, size_pt: int, bold: bool = False, color: RGBColor = DARK_GRAY,
               font_name: str = BODY_FONT) -> None:
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    _set_cjk_font(run, font_name)


def add_rect(slide, left_in, top_in, width_in, height_in, fill_color):
    """배경 사각형."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_header_bar(slide, title_text: str, slide_number: int) -> None:
    """상단 블루 헤더 바 + 타이틀."""
    # 얇은 블루 상단 액센트 라인
    add_rect(slide, 0, 0, SLIDE_W_IN, 0.25, ACCENT_BLUE)
    # 타이틀 텍스트
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4),
        Inches(SLIDE_W_IN - 1.5), Inches(0.85),
    )
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    _style_run(run, size_pt=TITLE_PT, bold=True, color=DARK_BLUE, font_name=TITLE_FONT)
    # 타이틀 아래 얇은 구분선
    add_rect(slide, 0.5, 1.3, SLIDE_W_IN - 1.0, 0.04, ACCENT_BLUE)
    # 우상단 슬라이드 번호 박스
    num_box = slide.shapes.add_textbox(
        Inches(SLIDE_W_IN - 1.2), Inches(0.4),
        Inches(0.9), Inches(0.5),
    )
    nt = num_box.text_frame
    nt.margin_left = Emu(0)
    nt.margin_right = Emu(0)
    np = nt.paragraphs[0]
    np.alignment = PP_ALIGN.RIGHT
    nr = np.add_run()
    nr.text = f"{slide_number} / 16"
    _style_run(nr, size_pt=FOOTER_PT + 2, bold=True, color=ACCENT_BLUE)


def add_footer(slide, slide_number: int) -> None:
    """하단 푸터 (프로젝트명 + 페이지)."""
    # 하단 얇은 라인
    add_rect(slide, 0.5, SLIDE_H_IN - 0.45, SLIDE_W_IN - 1.0, 0.02, ACCENT_BLUE)
    fb = slide.shapes.add_textbox(
        Inches(0.5), Inches(SLIDE_H_IN - 0.4),
        Inches(SLIDE_W_IN - 1.0), Inches(0.3),
    )
    ft = fb.text_frame
    ft.margin_left = Emu(0)
    ft.margin_right = Emu(0)
    fp = ft.paragraphs[0]
    fp.alignment = PP_ALIGN.LEFT
    fr = fp.add_run()
    fr.text = "SDACS \u2014 군집드론 공역통제 자동화 시스템"
    _style_run(fr, size_pt=FOOTER_PT, bold=False, color=GRAY)


def add_bullets(slide, bullets, left_in=0.7, top_in=1.6,
                width_in=7.0, height_in=5.0, size_pt=BODY_PT,
                line_spacing=1.3) -> None:
    """왼쪽 본문 bullet 목록."""
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)

    for i, item in enumerate(bullets):
        # item: str or (heading, subtext)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(8)

        if isinstance(item, tuple):
            head, sub = item
            bullet_run = p.add_run()
            bullet_run.text = "\u2022  "
            _style_run(bullet_run, size_pt=size_pt, bold=True, color=ACCENT_BLUE)
            head_run = p.add_run()
            head_run.text = head
            _style_run(head_run, size_pt=size_pt, bold=True, color=DARK_BLUE)
            if sub:
                sub_run = p.add_run()
                sub_run.text = "  " + sub
                _style_run(sub_run, size_pt=size_pt - 2, bold=False, color=DARK_GRAY)
        else:
            bullet_run = p.add_run()
            bullet_run.text = "\u2022  "
            _style_run(bullet_run, size_pt=size_pt, bold=True, color=ACCENT_BLUE)
            text_run = p.add_run()
            text_run.text = item
            _style_run(text_run, size_pt=size_pt, bold=False, color=DARK_GRAY)


def add_image(slide, img_key: str, left_in: float, top_in: float, max_w_in: float) -> None:
    """이미지 삽입 (max 폭 제한, 중앙 정렬된 위치)."""
    path = os.path.join(IMG_DIR, IMG[img_key])
    if not os.path.exists(path):
        print(f"[WARN] 이미지 없음: {path}")
        return
    slide.shapes.add_picture(
        path,
        Inches(left_in), Inches(top_in),
        width=Inches(max_w_in),
    )


def add_big_stat_box(slide, left_in: float, top_in: float,
                     width_in: float, height_in: float,
                     big_text: str, caption: str,
                     color: RGBColor = ACCENT_BLUE) -> None:
    """큰 숫자 강조 박스."""
    add_rect(slide, left_in, top_in, width_in, height_in, LIGHT_BLUE)
    # 얇은 왼쪽 액센트 바
    add_rect(slide, left_in, top_in, 0.15, height_in, color)

    tb = slide.shapes.add_textbox(
        Inches(left_in + 0.3), Inches(top_in + 0.15),
        Inches(width_in - 0.4), Inches(height_in - 0.3),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = big_text
    _style_run(r1, size_pt=36, bold=True, color=color)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = caption
    _style_run(r2, size_pt=14, bold=False, color=DARK_GRAY)


def new_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 1: 표지                                                 ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_1_title() -> None:
    slide = new_slide()

    # 전체 배경 블루
    add_rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, DARK_BLUE)
    # 하단 화이트 영역
    add_rect(slide, 0, SLIDE_H_IN - 2.0, SLIDE_W_IN, 2.0, WHITE)
    # 액센트 라인
    add_rect(slide, 0, SLIDE_H_IN - 2.05, SLIDE_W_IN, 0.1, ACCENT_BLUE)

    # 상단 레이블
    lbl = slide.shapes.add_textbox(
        Inches(1.0), Inches(0.9), Inches(11), Inches(0.5)
    )
    lt = lbl.text_frame
    lt.margin_left = Emu(0)
    lp = lt.paragraphs[0]
    lp.alignment = PP_ALIGN.LEFT
    lr = lp.add_run()
    lr.text = "캡스톤디자인 최종 발표 / 2026"
    _style_run(lr, size_pt=16, bold=False, color=RGBColor(0x9F, 0xC6, 0xE8))

    # 메인 타이틀
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = "SDACS"
    _style_run(r1, size_pt=72, bold=True, color=WHITE, font_name=TITLE_FONT)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "군집드론 공역통제 자동화 시스템"
    _style_run(r2, size_pt=32, bold=True, color=WHITE, font_name=TITLE_FONT)

    # 부제목
    sub_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(4.7), Inches(11), Inches(0.6)
    )
    st = sub_box.text_frame
    st.margin_left = Emu(0)
    sp = st.paragraphs[0]
    sp.alignment = PP_ALIGN.LEFT
    sr = sp.add_run()
    sr.text = "Swarm Drone Airspace Control System"
    _style_run(sr, size_pt=20, bold=False, color=RGBColor(0xB8, 0xD4, 0xEC))

    # 하단 저자/일자
    foot_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(SLIDE_H_IN - 1.6), Inches(11), Inches(1.3)
    )
    ft = foot_box.text_frame
    ft.margin_left = Emu(0)

    fp1 = ft.paragraphs[0]
    fp1.alignment = PP_ALIGN.LEFT
    fr1 = fp1.add_run()
    fr1.text = "발표자  |  장선우"
    _style_run(fr1, size_pt=18, bold=True, color=DARK_BLUE)

    fp2 = ft.add_paragraph()
    fp2.alignment = PP_ALIGN.LEFT
    fr2 = fp2.add_run()
    fr2.text = "지도교수  |  캡스톤디자인 프로젝트"
    _style_run(fr2, size_pt=14, bold=False, color=DARK_GRAY)

    fp3 = ft.add_paragraph()
    fp3.alignment = PP_ALIGN.LEFT
    fr3 = fp3.add_run()
    today_str = datetime.now().strftime("%Y년 %m월 %d일")
    fr3.text = f"발표일  |  {today_str}"
    _style_run(fr3, size_pt=14, bold=False, color=DARK_GRAY)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 2: 문제 정의                                            ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_2_problem() -> None:
    slide = new_slide()
    add_header_bar(slide, "1. 문제 정의 \u2014 하늘길이 막히고 있다", 2)

    add_bullets(slide, [
        ("드론 시장 폭발적 성장", "2030년 국내 드론 시장 4.5조 원 예상, 기체 수 100만 대 돌파"),
        ("배송·촬영·측량·방제 등 사용처 급증", "도심 저고도(150 m 이하) 공역이 급속히 혼잡해짐"),
        ("공역 용량 한계 도달", "서울·광주 등 주요 도심 공역은 이미 포화 직전"),
        ("충돌·추락 사고 리스크 증가", "현재 제도로는 100대 이상 동시 비행 안전 보장 불가"),
    ], left_in=0.7, top_in=1.6, width_in=7.2, height_in=5.2)

    add_image(slide, "g1", left_in=8.0, top_in=1.7, max_w_in=4.9)

    add_footer(slide, 2)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 3: 기존 방식 한계                                       ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_3_gap() -> None:
    slide = new_slide()
    add_header_bar(slide, "2. 기존 방식의 한계", 3)

    add_bullets(slide, [
        ("수동 관제에 의존", "관제사 1명이 드론 수십 대를 실시간 감시 \u2014 휴먼 에러 불가피"),
        ("단일 드론 중심 제어", "드론 한 대씩 명령 전달 \u2014 군집 단위 조율 불가능"),
        ("지상 레이더 한계", "도심 빌딩이 저고도 공역 67%를 감시 사각지대로 만듦"),
        ("규정·소프트웨어 미비", "100대 이상 동시 비행을 관제할 표준 체계가 없음"),
    ], left_in=0.7, top_in=1.6, width_in=7.2, height_in=5.2)

    add_image(slide, "g2", left_in=8.0, top_in=1.7, max_w_in=4.9)

    add_footer(slide, 3)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 4: SDACS 해결책 개요                                    ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_4_solution() -> None:
    slide = new_slide()
    add_header_bar(slide, "3. SDACS 해결책 \u2014 4계층 자동화 구조", 4)

    add_bullets(slide, [
        ("공역통제의 완전 자동화", "사람이 아닌 시스템이 충돌·진입·퇴각을 실시간 판단"),
        ("4계층 역할 분담 구조", "드론 → 공역 제어기 → 시뮬레이터 → UI 대시보드"),
        ("게임 AI 연구를 현실로 이식", "StarCraft II 군집 제어 알고리즘을 실제 드론 관제에 적용"),
        ("검증된 오픈소스 기반", "SimPy · Dash · WebGPU \u2014 재현 가능 + 확장 가능"),
    ], left_in=0.7, top_in=1.6, width_in=7.2, height_in=5.2)

    add_image(slide, "g4", left_in=8.0, top_in=1.7, max_w_in=4.9)

    add_footer(slide, 4)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 5: Layer 1 드론 에이전트                                ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_5_layer1() -> None:
    slide = new_slide()
    add_header_bar(slide, "4. Layer 1 \u2014 드론 에이전트 (10 Hz)", 5)

    add_bullets(slide, [
        ("SimPy 기반 이산 이벤트 시뮬레이션", "각 드론이 독립적인 Python 프로세스로 동작"),
        ("10 Hz 업데이트 루프", "100 ms마다 위치·속도·배터리·풍속 영향 갱신"),
        ("자율 충돌 회피", "APF 기반 개별 드론 회피 반응 (100 ms 응답)"),
        ("센서 모델링", "GPS 노이즈 · 배터리 드레인 · 풍속 외란까지 포함"),
        ("확장성", "50 \u2192 500대까지 선형 확장 검증 완료"),
    ], left_in=0.7, top_in=1.6, width_in=7.2, height_in=5.2)

    add_image(slide, "g3", left_in=8.0, top_in=1.7, max_w_in=4.9)

    add_footer(slide, 5)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 6: Layer 2 공역 제어기                                  ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_6_layer2() -> None:
    slide = new_slide()
    add_header_bar(slide, "5. Layer 2 \u2014 공역 제어기 (1 Hz, APF)", 6)

    add_bullets(slide, [
        ("AirspaceController 1 Hz 주기", "매초 전체 군집의 잠재적 충돌 분석 + 경로 권고"),
        ("Artificial Potential Field 기반", "목표지점은 인력(attract), 장애물·타 드론은 척력(repel)"),
        ("5겹 안전망 적용", "권고 \u2192 속도 조절 \u2192 호버링 \u2192 회피기동 \u2192 강제 착륙"),
        ("우선순위 자동 조정", "배터리 잔량·임무 중요도·위험도에 따라 5단계 동적 변환"),
        ("실시간 충돌 해결", "평균 해결률 97.8% (63개 시나리오 평균)"),
    ], left_in=0.7, top_in=1.6, width_in=7.2, height_in=5.2)

    add_image(slide, "g5", left_in=8.0, top_in=1.7, max_w_in=4.9)

    add_footer(slide, 6)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 7: Layer 3 시뮬레이터/몬테카를로                         ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_7_layer3() -> None:
    slide = new_slide()
    add_header_bar(slide, "6. Layer 3 \u2014 시뮬레이터 & Monte Carlo", 7)

    add_bullets(slide, [
        ("SwarmSimulator 단일 엔진", "SimPy 이벤트 루프 + 풍속 모델 + 충돌 감지 통합"),
        ("WindModel 환경 외란", "Perlin 노이즈 + 가우시안 외란으로 실측 풍장 근사"),
        ("Monte Carlo 파라미터 스윕", "38,400 runs \u2014 드론 수·풍속·밀도 전 조합 검증"),
        ("재현성 보장", "np.random.default_rng(seed) 기반, 동일 입력 \u2192 동일 결과"),
        ("통계적 성능 지표", "충돌률·해결률·응답 지연·처리량을 자동 집계"),
    ], left_in=0.7, top_in=1.6, width_in=7.2, height_in=5.2)

    add_image(slide, "g7", left_in=8.0, top_in=1.7, max_w_in=4.9)

    add_footer(slide, 7)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 8: APF 알고리즘                                         ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_8_apf() -> None:
    slide = new_slide()
    add_header_bar(slide, "7. APF (Artificial Potential Field) 알고리즘", 8)

    add_bullets(slide, [
        ("전체 공역을 힘의 장(field)으로 모델링", "목표지점 = 끌어당기는 자석, 장애물 = 밀어내는 자석"),
        ("드론은 힘의 합력 방향으로 이동", "각 프레임마다 그래디언트 계산 \u2192 속도 벡터 갱신"),
        ("실시간 계산 가능", "벡터 연산 기반 \u2014 500대까지 1초 내 완료"),
        ("국소 최소점 탈출 로직", "목표 도달 실패 시 임의 섭동(perturbation) 주입"),
        ("기존 경로계획 대비 장점", "재계산 비용이 낮아 동적 환경에 강함"),
    ], left_in=0.7, top_in=1.6, width_in=7.2, height_in=5.2)

    add_image(slide, "g6", left_in=8.0, top_in=1.7, max_w_in=4.9)

    add_footer(slide, 8)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 9: 강풍 모드 자동 전환                                   ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_9_wind() -> None:
    slide = new_slide()
    add_header_bar(slide, "8. 강풍 모드 자동 전환 (풍속 > 10 m/s)", 9)

    add_bullets(slide, [
        ("기상 적응형 파라미터", "WindModel 측정값이 10 m/s 초과 시 자동으로 강풍 모드"),
        ("APF_PARAMS_WINDY 프리셋", "척력 계수 1.8배 \u00b7 회피 반경 1.5배 \u00b7 속도 상한 0.7배"),
        ("안전 vs 처리량 균형", "강풍 중에도 해결률 95% 이상 유지"),
        ("수동 개입 불필요", "운영자 조작 없이 실시간 자동 전환"),
        ("광주광역시 실측 풍장 검증", "평균 돌풍 12 m/s 환경에서 충돌 0건 달성"),
    ], left_in=0.7, top_in=1.6, width_in=7.2, height_in=5.2)

    add_image(slide, "g9", left_in=8.0, top_in=1.7, max_w_in=4.9)

    add_footer(slide, 9)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 10: WebGPU Compute 최적화                                ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_10_gpu() -> None:
    slide = new_slide()
    add_header_bar(slide, "9. WebGPU Compute Shader 가속", 10)

    add_bullets(slide, [
        ("CPU 병목을 GPU로 오프로드", "APF 그래디언트 계산을 WebGPU 셰이더에서 병렬 수행"),
        ("병렬 처리 효과", "500대 동시 계산 \u2014 CPU 대비 약 12배 처리량 향상"),
        ("브라우저 기반 3D 시각화", "Chrome · Edge · Firefox에서 별도 설치 없이 구동"),
        ("실시간 대시보드", "60 FPS 이상 \u2014 관제사 화면에서 부드러운 군집 시각화"),
        ("확장성 확보", "GPU 연산 자원이 남는 한 드론 수 수천 대까지 대응 가능"),
    ], left_in=0.7, top_in=1.6, width_in=7.2, height_in=5.2)

    add_image(slide, "g8", left_in=8.0, top_in=1.7, max_w_in=4.9)

    add_footer(slide, 10)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 11: 7대 광역시 시나리오                                  ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_11_cities() -> None:
    slide = new_slide()
    add_header_bar(slide, "10. 7대 광역시 실증 시나리오", 11)

    add_bullets(slide, [
        ("서울·부산·대구·인천·광주·대전·울산", "7개 대도시 지형/풍장/빌딩 데이터를 시뮬레이션에 반영"),
        ("카테고리 63종 시나리오", "고밀도·강풍·불법 진입·긴급 복귀 등 실제 발생 가능 상황"),
        ("도시별 평균 성공률 96.9%", "최저 서울 95.4% \u2014 최고 대전 98.7%"),
        ("광주광역시 집중 검증", "무진대로 2 km 회랑 20대 군집 1시간 연속 비행 \u2014 충돌 0건"),
        ("실증 준비 완료", "시뮬레이션 \u2192 실제 비행 테스트로 이관 가능한 단계"),
    ], left_in=0.7, top_in=1.6, width_in=7.2, height_in=5.2)

    add_image(slide, "g10", left_in=8.0, top_in=1.7, max_w_in=4.9)

    add_footer(slide, 11)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 12: 핵심 성과 97.8%                                     ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_12_success() -> None:
    slide = new_slide()
    add_header_bar(slide, "11. 핵심 성과 \u2014 충돌 해결률 97.8%", 12)

    # 좌측: 강조 숫자 박스 3개
    add_big_stat_box(slide, 0.7, 1.7, 4.2, 1.5,
                     "97.8 %", "충돌 해결률 (63 시나리오 평균)",
                     color=ACCENT_BLUE)
    add_big_stat_box(slide, 0.7, 3.4, 4.2, 1.5,
                     "38,400", "Monte Carlo 실험 횟수",
                     color=DARK_BLUE)
    add_big_stat_box(slide, 0.7, 5.1, 4.2, 1.5,
                     "0 건", "광주 2 km 회랑 1시간 충돌 건수",
                     color=RED_ACCENT)

    # 우측: 설명 bullets
    add_bullets(slide, [
        ("해결률 산식", "1 - collisions / (conflicts + collisions)"),
        ("500대 대규모 군집까지 검증", "성공률 선형 유지 \u2014 드론 수에 덜 민감"),
        ("통계적 유의성 확보", "38,400 runs \u2192 95% 신뢰구간 \u00b1 0.4%p"),
        ("재현 가능성", "seed 고정 \u2192 누구나 동일 결과 재현"),
    ], left_in=5.5, top_in=1.7, width_in=7.5, height_in=5.0, size_pt=17)

    add_footer(slide, 12)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 13: 응답 시간 + 처리량                                   ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_13_latency() -> None:
    slide = new_slide()
    add_header_bar(slide, "12. 응답 시간 ~ 1.2 초 · 처리량 1.05배", 13)

    add_big_stat_box(slide, 0.7, 1.7, 4.2, 1.5,
                     "1.2 초", "탐지 \u2192 권고 평균 응답 시간",
                     color=ACCENT_BLUE)
    add_big_stat_box(slide, 0.7, 3.4, 4.2, 1.5,
                     "1.05 \u00d7", "동일 공역 대비 처리량 증가",
                     color=DARK_BLUE)
    add_big_stat_box(slide, 0.7, 5.1, 4.2, 1.5,
                     "< 100 ms", "드론 개별 APF 반응 시간",
                     color=RED_ACCENT)

    add_bullets(slide, [
        ("5단계 자동 파이프라인", "탐지 \u2192 식별 \u2192 판단 \u2192 권고 \u2192 실행 (전체 < 1.5초)"),
        ("동일 공역 처리량 증가", "관제 자동화로 동일 영역 내 +5% 비행 밀도 허용"),
        ("관제사 업무량 감소", "권고는 자동 생성 \u2192 사람은 예외 상황만 처리"),
        ("스케줄 신뢰도 향상", "응답 지연 표준편차 0.3초 \u2014 예측 가능한 관제"),
    ], left_in=5.5, top_in=1.7, width_in=7.5, height_in=5.0, size_pt=17)

    add_footer(slide, 13)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 14: 3D 시각화 스크린샷                                   ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_14_viz() -> None:
    slide = new_slide()
    add_header_bar(slide, "13. 3D 시뮬레이션 시각화", 14)

    add_bullets(slide, [
        ("Dash + WebGPU 3D 대시보드", "브라우저에서 실시간 군집 비행 확인"),
        ("관제사용 인터페이스", "드론 ID · 속도 · 배터리 · 경로 · 경보를 한 화면에"),
        ("카메라 조작·시간 슬라이더", "과거 사건 되감기·시점 전환으로 원인 분석 가능"),
        ("3D 도시 모델 연동", "서울·광주 등 실제 빌딩 데이터 로드"),
    ], left_in=0.7, top_in=1.55, width_in=7.0, height_in=5.5)

    add_image(slide, "g13", left_in=7.9, top_in=1.7, max_w_in=5.0)

    add_footer(slide, 14)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 15: 향후 계획                                           ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_15_future() -> None:
    slide = new_slide()
    add_header_bar(slide, "14. 향후 계획 \u2014 특허 · 실증 · 확장", 15)

    add_bullets(slide, [
        ("특허 출원 5건 준비", "APF 강풍 자동 전환 / 5겹 안전망 / WebGPU 관제 UI 등"),
        ("광주광역시 실증 실험", "2 km 회랑 20대 실기체 시범 비행 \u2014 안전 인증 절차 병행"),
        ("규정·표준화 연계", "국토부·항공안전기술원과 공역통제 자동화 가이드라인 논의"),
        ("산업 확장 로드맵", "배송·방제·재난구조·공공안전 등 다영역 적용"),
        ("오픈소스 공개", "GitHub 공개 \u2013 후속 연구진·스타트업과 협력"),
    ], left_in=0.7, top_in=1.6, width_in=7.2, height_in=5.2)

    add_image(slide, "g11", left_in=8.0, top_in=1.7, max_w_in=4.9)

    add_footer(slide, 15)


# ╔═══════════════════════════════════════════════════════════════╗
# ║ Slide 16: 결론 + 감사합니다                                   ║
# ╚═══════════════════════════════════════════════════════════════╝
def build_slide_16_conclusion() -> None:
    slide = new_slide()

    # 배경: 좌측 블루 / 우측 화이트
    add_rect(slide, 0, 0, 5.5, SLIDE_H_IN, DARK_BLUE)
    add_rect(slide, 5.5, 0, SLIDE_W_IN - 5.5, SLIDE_H_IN, WHITE)
    add_rect(slide, 5.5, 0, 0.08, SLIDE_H_IN, ACCENT_BLUE)

    # 좌측: "감사합니다"
    lt = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(5.0), Inches(2.5)
    )
    lf = lt.text_frame
    lf.word_wrap = True

    lp1 = lf.paragraphs[0]
    lp1.alignment = PP_ALIGN.LEFT
    lr1 = lp1.add_run()
    lr1.text = "감사합니다"
    _style_run(lr1, size_pt=48, bold=True, color=WHITE, font_name=TITLE_FONT)

    lp2 = lf.add_paragraph()
    lp2.alignment = PP_ALIGN.LEFT
    lr2 = lp2.add_run()
    lr2.text = "Thank you"
    _style_run(lr2, size_pt=22, bold=False, color=RGBColor(0x9F, 0xC6, 0xE8))

    # 좌측 하단: 연락처
    ct = slide.shapes.add_textbox(
        Inches(0.5), Inches(SLIDE_H_IN - 1.5), Inches(5.0), Inches(1.2)
    )
    cf = ct.text_frame
    cp = cf.paragraphs[0]
    cp.alignment = PP_ALIGN.LEFT
    cr = cp.add_run()
    cr.text = "SDACS Project  |  2026"
    _style_run(cr, size_pt=14, bold=False, color=RGBColor(0x9F, 0xC6, 0xE8))

    # 우측: 결론 요약
    rt = slide.shapes.add_textbox(
        Inches(6.0), Inches(1.2), Inches(7.0), Inches(5.5)
    )
    rf = rt.text_frame
    rf.word_wrap = True

    rp_title = rf.paragraphs[0]
    rp_title.alignment = PP_ALIGN.LEFT
    rr_title = rp_title.add_run()
    rr_title.text = "15. 결론"
    _style_run(rr_title, size_pt=24, bold=True, color=DARK_BLUE)

    conclusions = [
        ("자동화된 공역통제는 실현 가능", "시뮬레이션·검증·시각화까지 end-to-end 구축"),
        ("안전성 입증", "충돌 해결률 97.8% · 500대 군집 · 강풍까지 견딤"),
        ("확장 가능한 구조", "SimPy + WebGPU 기반 오픈소스 아키텍처"),
        ("학계·산업·규제 연결 기술", "특허 · 실증 · 표준화 3축 동시 추진"),
    ]
    for head, sub in conclusions:
        p = rf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.3
        p.space_before = Pt(10)

        b = p.add_run()
        b.text = "\u2022  "
        _style_run(b, size_pt=16, bold=True, color=ACCENT_BLUE)
        h = p.add_run()
        h.text = head
        _style_run(h, size_pt=16, bold=True, color=DARK_BLUE)
        s = p.add_run()
        s.text = "  " + sub
        _style_run(s, size_pt=14, bold=False, color=DARK_GRAY)


# ── 빌드 실행 ────────────────────────────────────────
def build() -> None:
    build_slide_1_title()
    build_slide_2_problem()
    build_slide_3_gap()
    build_slide_4_solution()
    build_slide_5_layer1()
    build_slide_6_layer2()
    build_slide_7_layer3()
    build_slide_8_apf()
    build_slide_9_wind()
    build_slide_10_gpu()
    build_slide_11_cities()
    build_slide_12_success()
    build_slide_13_latency()
    build_slide_14_viz()
    build_slide_15_future()
    build_slide_16_conclusion()

    prs.save(OUT_PATH)
    size_kb = os.path.getsize(OUT_PATH) / 1024.0
    print(f"[OK] saved: {OUT_PATH}")
    print(f"[OK] slides: {len(prs.slides)}")
    print(f"[OK] size: {size_kb:.1f} KB")


if __name__ == "__main__":
    build()
